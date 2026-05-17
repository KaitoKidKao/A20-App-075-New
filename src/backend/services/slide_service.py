import logging
import json
import copy
import random
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List
from pptx import Presentation
from pptx.text.text import _Paragraph
from sqlmodel import Session, select
from ..models import Lesson, ContentMetadata
from .ai_service import AIService
from .storage_service import ensure_local, upload_to_s3
from .. import config

logger = logging.getLogger(__name__)


class SlideService:
    TEMPLATES_DIR = Path("src/backend/assets/slide_templates")
    OUTPUT_DIR = Path("data/uploads/slides")
    S3_SLIDES_PREFIX = "AI_Result/slides"

    @staticmethod
    async def generate_slides(
        session: Session, lesson_id: str, template_id: str, num_slides: int = 10
    ) -> Dict[str, Any]:
        # Parse UUID — fixes string vs UUID column mismatch in PostgreSQL
        try:
            lesson_uuid = uuid.UUID(str(lesson_id))
        except ValueError:
            raise ValueError(f"Invalid lesson id: {lesson_id}")

        # 1. Resolve template
        template_path = SlideService.TEMPLATES_DIR / f"{template_id}.pptx"
        if not template_path.exists():
            templates = sorted(SlideService.TEMPLATES_DIR.glob("*.pptx"))
            if not templates:
                raise ValueError(f"Không tìm thấy template tại {SlideService.TEMPLATES_DIR}")
            template_path = templates[0]
            template_id = template_path.stem

        # 2. Verify lesson
        lesson = session.get(Lesson, lesson_uuid)
        if not lesson:
            raise ValueError("Không tìm thấy bài học.")

        # 3. Load transcript — local first, then S3 fallback
        transcript_path = AIService.TRANSCRIPT_DIR / f"{lesson_id}.json"
        s3_key = f"Transcripts/{lesson_id}.json"
        if not transcript_path.exists():
            logger.info(f"Transcript not local, trying S3: {s3_key}")
            ensure_local(s3_key, transcript_path)

        if not transcript_path.exists():
            raise ValueError("Chưa có transcript để tạo slide.")

        with open(transcript_path, "r", encoding="utf-8") as f:
            transcript_data = json.load(f)

        # 4. LLM generates slide content
        logger.info(f"🧠 LLM generating {num_slides} slides for lesson {lesson_id}...")
        slides_content = await SlideService._call_llm_for_slides(transcript_data, num_slides)

        # Pad if LLM returned fewer slides than requested
        while len(slides_content) < num_slides:
            if len(slides_content) >= 2:
                insert_pos = len(slides_content) - 1  # before last slide
                prev_title = slides_content[max(insert_pos - 1, 0)]["title"]
                slides_content.insert(insert_pos, {
                    "title": f"{prev_title} (Tiếp theo)",
                    "content": [
                        "Mở rộng thêm các ý chi tiết của phần trước...",
                        "Ví dụ minh họa và ứng dụng thực tế.",
                    ],
                })
            else:
                slides_content.append({
                    "title": "Nội dung bổ sung",
                    "content": ["Chi tiết bài giảng đang được cập nhật..."],
                })

        # Trim if LLM returned more slides than requested
        if len(slides_content) > num_slides:
            last_slide = slides_content[-1]
            slides_content = slides_content[: num_slides - 1]
            slides_content.append(last_slide)

        # 5. Assemble PPTX
        SlideService.OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
        output_filename = f"{lesson_id}_{uuid.uuid4().hex[:8]}.pptx"
        output_path = SlideService.OUTPUT_DIR / output_filename

        try:
            SlideService._assemble_pptx(template_path, slides_content, output_path)
        except Exception as e:
            logger.error(f"❌ PPTX assembly error: {e}")
            raise RuntimeError(f"Lỗi khi lắp ghép slide: {e}")

        # 6. Upload to S3 (AI_Result/slides/)
        slide_s3_key = f"{SlideService.S3_SLIDES_PREFIX}/{output_filename}"
        upload_to_s3(output_path, slide_s3_key)

        # 7. Persist metadata to DB
        metadata = session.exec(
            select(ContentMetadata).where(ContentMetadata.lesson_id == lesson_uuid)
        ).first()
        if not metadata:
            metadata = ContentMetadata(lesson_id=lesson_uuid, ai_analysis={})
            session.add(metadata)

        new_analysis = dict(metadata.ai_analysis or {})
        new_analysis["slides"] = {
            "template_id": template_id,
            "filename": output_filename,
            "s3_key": slide_s3_key,
            "generated_at": datetime.now(timezone.utc).isoformat(),
            "num_slides": len(slides_content),
        }
        metadata.ai_analysis = new_analysis
        session.add(metadata)
        session.commit()

        return {
            "status": "completed",
            "filename": output_filename,
            "s3_key": slide_s3_key,
            "num_slides": len(slides_content),
        }

    @staticmethod
    def _assemble_pptx(template_path: Path, slides_content: List[Dict], output_path: Path):
        new_prs = Presentation(template_path)
        num_template_slides = len(new_prs.slides)
        num_needed = len(slides_content)

        # Clone slides if template has fewer than needed
        if num_needed > num_template_slides and num_template_slides > 1:
            closing_slide_idx = num_template_slides - 1
            content_indices = list(range(1, num_template_slides - 1)) or [0]

            for _ in range(num_needed - num_template_slides):
                source_idx = random.choice(content_indices)
                source_slide = new_prs.slides[source_idx]
                new_slide = new_prs.slides.add_slide(source_slide.slide_layout)
                try:
                    if source_slide._element.bg is not None:
                        new_slide._element.insert(0, copy.deepcopy(source_slide._element.bg))
                except Exception:
                    pass
                for shape in source_slide.shapes:
                    new_el = copy.deepcopy(shape.element)
                    new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
                for rel in source_slide.part.rels.values():
                    if "image" in rel.reltype:
                        try:
                            new_slide.part.rels.get_or_add(rel.reltype, rel._target, rel.rId)
                        except Exception:
                            pass

            # Keep closing slide last
            xml_slides = new_prs.slides._sldIdLst
            closing_id = xml_slides[closing_slide_idx]
            xml_slides.remove(closing_id)
            xml_slides.append(closing_id)

        # Delete excess slides (keep last slide as closing)
        elif num_needed < num_template_slides and num_needed > 1:
            num_to_delete = num_template_slides - num_needed
            for _ in range(num_to_delete):
                idx = num_needed - 1
                rId = new_prs.slides._sldIdLst[idx].rId
                new_prs.part.drop_rel(rId)
                del new_prs.slides._sldIdLst[idx]

        final_slides = list(new_prs.slides)

        for i, content in enumerate(slides_content):
            if i >= len(final_slides):
                break
            slide = final_slides[i]

            text_shapes = [s for s in slide.shapes if hasattr(s, "text_frame") and s.text.strip()]
            text_shapes.sort(key=lambda s: s.top)

            title_shape = text_shapes[0] if text_shapes else slide.shapes.title
            body_shape = None
            if len(text_shapes) > 1:
                others = sorted(text_shapes[1:], key=lambda s: s.width * s.height, reverse=True)
                body_shape = others[0]
            else:
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx in (1, 2, 3):
                        body_shape = shape
                        break

            # Clear unused text boxes
            for s in text_shapes:
                if s is not title_shape and s is not body_shape:
                    try:
                        s.text = ""
                    except Exception:
                        pass

            # Fill title
            if title_shape:
                tf = title_shape.text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    p = tf.paragraphs[0]
                    p.runs[0].text = content.get("title", "")
                    for run in p.runs[1:]:
                        run._r.getparent().remove(run._r)
                else:
                    tf.paragraphs[0].text = content.get("title", "")

            # Fill body
            if body_shape:
                tf = body_shape.text_frame
                points = content.get("content", [])
                if not tf.paragraphs:
                    tf.text = ""

                template_p_xml = tf.paragraphs[0]._p
                parent_elt = template_p_xml.getparent()

                for p in list(tf.paragraphs):
                    p._p.getparent().remove(p._p)

                items = points if isinstance(points, list) else [str(points)]
                for p_text in items:
                    new_p_xml = copy.deepcopy(template_p_xml)
                    parent_elt.append(new_p_xml)
                    new_p = _Paragraph(new_p_xml, tf)
                    full_text = f"• {p_text}"
                    if new_p.runs:
                        new_p.runs[0].text = full_text
                        for run in new_p.runs[1:]:
                            run._r.getparent().remove(run._r)
                    else:
                        new_p.text = full_text

        new_prs.save(output_path)

    @staticmethod
    async def _call_llm_for_slides(transcript_data: dict, num_slides: int) -> List[Dict]:
        from openai import AsyncOpenAI

        if not config.OPENAI_API_KEY:
            return [{"title": "Lỗi", "content": ["Chưa cấu hình API Key."]}]

        client = AsyncOpenAI(api_key=config.OPENAI_API_KEY)

        full_text = " ".join([s["text"] for s in transcript_data.get("segments", [])])
        truncated_text = full_text[:8000]

        prompt = f"""
Bạn là một chuyên gia soạn thảo bài giảng. Hãy soạn nội dung cho một bộ Slide gồm CHÍNH XÁC {num_slides} trang dựa trên nội dung bài giảng dưới đây.

Yêu cầu:
1. Trả về DUY NHẤT một mảng JSON có đúng {num_slides} đối tượng. Mỗi đối tượng gồm:
   - "title": Tiêu đề slide.
   - "content": Danh sách các ý chính (mảng các chuỗi, tối đa 5 ý).
2. Slide 1 phải là slide tiêu đề tổng quát.
3. Slide cuối cùng (Slide thứ {num_slides}) phải là slide kết luận và cảm ơn.
4. Nội dung phải cô đọng, súc tích, phù hợp để trình chiếu.
5. TẤT CẢ NỘI DUNG PHẢI BẰNG TIẾNG VIỆT.

Nội dung bài giảng:
{truncated_text}
"""

        try:
            response = await client.chat.completions.create(
                model=config.DEFAULT_MODEL,
                messages=[
                    {"role": "system", "content": "Bạn là chuyên gia giáo dục. Trả về JSON."},
                    {"role": "user", "content": prompt},
                ],
                response_format={"type": "json_object"},
            )
            data = json.loads(response.choices[0].message.content)
            if isinstance(data, list):
                return data
            # LLM may wrap in {"slides": [...]} or {"items": [...]}
            for key in ("slides", "items"):
                if key in data and isinstance(data[key], list):
                    return data[key]
            # Fallback: take first list value
            for val in data.values():
                if isinstance(val, list):
                    return val
            return [{"title": "Lỗi", "content": ["Dữ liệu không hợp lệ từ LLM."]}]
        except Exception as e:
            logger.error(f"❌ LLM slide error: {e}")
            return [{"title": "Lỗi", "content": ["Không thể sinh nội dung slide."]}]
